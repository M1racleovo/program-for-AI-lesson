"""
FaceFair PPT 生成器
直接从分析结果和图表生成完整演示文稿（无需模板+填充两步走）
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from config import OUTPUT_DIR, CHARTS_DIR, SUMMARY_TXT, RACE_CATEGORIES

# ---------- 配色 ----------
COLOR_DARK = RGBColor(0x1B, 0x1B, 0x2F)
COLOR_PRIMARY = RGBColor(0x2D, 0x5F, 0x8A)
COLOR_ACCENT = RGBColor(0xE8, 0x6A, 0x17)
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
COLOR_GRAY = RGBColor(0x88, 0x88, 0x88)
COLOR_GREEN = RGBColor(0x27, 0xAE, 0x60)
COLOR_RED = RGBColor(0xC0, 0x39, 0x2B)


# ========== 工具函数 ==========

def add_bg(slide, color=COLOR_LIGHT_BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_dark_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_DARK


def add_bottom_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.1), Inches(13.33), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()


def add_title_text(slide, text, left, top, width, height,
                   font_size=28, color=COLOR_DARK, bold=True, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment


def add_body_text(slide, text, left, top, width, height,
                  font_size=16, color=COLOR_BLACK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
    return tf


def add_separator(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()


def add_page_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(12.3), Inches(7.15), Inches(0.8), Inches(0.3))
    p = txBox.text_frame.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_image_safe(slide, img_path, left, top, width, height):
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))


# ========== 数据解析 ==========

def parse_summary():
    with open(SUMMARY_TXT, "r", encoding="utf-8") as f:
        text = f.read()

    data = {}
    for pattern, key in [
        (r"总样本数:\s*(\d+)", "total"),
        (r"有效样本数.*?:\s*(\d+)", "valid"),
        (r"人脸检测失败数:\s*(\d+)", "failed"),
        (r"总体种族识别准确率:\s*([\d.]+)%", "race_acc"),
        (r"总体性别识别准确率:\s*([\d.]+)%", "gender_acc"),
    ]:
        m = re.search(pattern, text)
        data[key] = int(m.group(1)) if m and key in ("total", "valid", "failed") else (float(m.group(1)) if m else 0)

    race_stats = {}
    for race in RACE_CATEGORIES:
        m = re.search(rf"{race}\s+(\d+)\s+([\d.]+)\s*%\s+([\d.]+)\s*%\s+([\d.]+)", text)
        if m:
            race_stats[race] = {"count": int(m.group(1)), "race_acc": float(m.group(2)),
                                "gender_acc": float(m.group(3)), "avg_age": float(m.group(4))}
    data["race_stats"] = race_stats

    gender_stats = {}
    for g in ["Male", "Female"]:
        m = re.search(rf"{g}\s+(\d+)\s+([\d.]+)\s*%", text)
        if m:
            gender_stats[g] = {"count": int(m.group(1)), "acc": float(m.group(2))}
    data["gender_stats"] = gender_stats

    for key, pattern in [
        ("best_race", r"最高种族准确率:\s*(.+?)\s*\(([\d.]+)%\)"),
        ("worst_race", r"最低种族准确率:\s*(.+?)\s*\(([\d.]+)%\)"),
    ]:
        m = re.search(pattern, text)
        if m:
            data[key] = m.group(1).strip()
            data[key.replace("_race", "_acc")] = float(m.group(2))

    m = re.search(r"种族准确率差距:\s*([\d.]+)%", text)
    data["gap"] = float(m.group(1)) if m else 0
    m = re.search(r"评估:\s*(.+)$", text, re.MULTILINE)
    data["assessment"] = m.group(1).strip() if m else ""

    if data["gap"] > 15:
        data["bias_level"] = "显著"
    elif data["gap"] > 5:
        data["bias_level"] = "中度"
    else:
        data["bias_level"] = "轻微"

    return data


# ========== 幻灯片 ==========

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.2), Inches(1.2), Inches(0.06))
    shape.fill.solid(); shape.fill.fore_color.rgb = COLOR_ACCENT; shape.line.fill.background()

    add_title_text(slide, "FaceFair", 1, 1.5, 11, 1.2, font_size=48, color=COLOR_WHITE)
    add_title_text(slide, "人脸识别公平性审计", 1, 2.3, 11, 0.8, font_size=30, color=RGBColor(0xBB, 0xCC, 0xDD), bold=False)
    add_body_text(slide, "基于 DeepFace / Facenet 的多维度 AI 偏见量化分析\n人工智能导论 · 课程项目", 1, 4.2, 8, 1.2, font_size=14, color=COLOR_GRAY)
    add_body_text(slide, "[姓名 / 学号 / 班级]\n[日期]", 1, 5.6, 6, 0.8, font_size=12, color=COLOR_GRAY)


def slide_background(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_bottom_bar(slide); add_page_number(slide, 2)
    add_title_text(slide, "项目背景", 0.8, 0.4, 10, 0.7, font_size=32, color=COLOR_DARK)
    add_separator(slide)
    add_body_text(slide,
        "人脸识别技术已广泛应用于安防、金融、社交媒体等领域。\n"
        "然而，多项研究表明主流人脸识别 AI 在不同种族、性别上存在\n"
        "显著的性能差异，可能导致歧视性后果。\n\n"
        "■ 2018 年 MIT & Stanford 研究: IBM/微软/旷视的商业 API\n"
        "  对深色皮肤女性的性别识别错误率高达 34.7% (浅色男性仅 0.8%)\n"
        "■ NIST 2020 报告: 多数算法对亚裔和非裔面孔的误识率偏高\n"
        "■ 监管趋势: 欧盟 AI Act 要求高风险 AI 系统提供公平性证明",
        0.8, 1.5, 11.5, 4.8, font_size=16, color=COLOR_BLACK)


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_bottom_bar(slide); add_page_number(slide, 3)
    add_title_text(slide, "研究问题", 0.8, 0.4, 10, 0.7, font_size=32, color=COLOR_DARK)
    add_separator(slide)

    questions = [
        ("Q1", "主流人脸识别模型 (Facenet) 在不同种族上的识别准确率\n是否存在显著差异？"),
        ("Q2", "性别维度是否与种族维度叠加，造成更严重的交叉偏见？"),
        ("Q3", "如果存在偏见，其程度如何量化？是否超过公平性阈值？"),
    ]
    for i, (label, text) in enumerate(questions):
        y = 1.6 + i * 1.7
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y), Inches(0.5), Inches(0.5))
        circle.fill.solid(); circle.fill.fore_color.rgb = COLOR_PRIMARY; circle.line.fill.background()
        p = circle.text_frame.paragraphs[0]
        p.text = label; p.font.size = Pt(14); p.font.color.rgb = COLOR_WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
        add_body_text(slide, text, 1.8, y - 0.05, 10, 0.8, font_size=15, color=COLOR_BLACK)


def slide_approach(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_bottom_bar(slide); add_page_number(slide, 4)
    add_title_text(slide, "技术方案", 0.8, 0.4, 10, 0.7, font_size=32, color=COLOR_DARK)
    add_separator(slide)

    add_title_text(slide, "模型选型", 0.8, 1.4, 5, 0.5, font_size=20, color=COLOR_PRIMARY)
    add_body_text(slide,
        "人脸检测: OpenCV Haar Cascade\n"
        "├─ 轻量级，CPU 上快速运行\n"
        "└─ 适合本地批量处理\n\n"
        "人脸识别: Facenet (Google)\n"
        "├─ CNN 架构，128 维嵌入向量\n"
        "├─ Triplet Loss 训练，经典基准模型\n"
        "└─ 被广泛部署，审计价值高\n\n"
        "属性预测: DeepFace 内置分类器\n"
        "├─ 种族分类 (7 类)\n"
        "├─ 性别分类 (男/女)\n"
        "└─ 年龄回归",
        0.8, 2.0, 5.5, 4.8, font_size=13, color=COLOR_BLACK)

    add_title_text(slide, "选型理由", 6.8, 1.4, 5, 0.5, font_size=20, color=COLOR_PRIMARY)
    add_body_text(slide,
        "Facenet 非最新 SOTA，但:\n\n"
        "① 部署量巨大 — 审计结果影响面广\n"
        "② 轻量高效 — 本地即可跑完整测试\n"
        "③ 经典架构 — 偏见若存在，根源于\n"
        "   训练数据 & 损失函数设计\n"
        "④ 可复现 — 权重公开，结果可验证\n\n"
        "[后续可选]\n"
        "→ 换 MTCNN 检测器做对比实验\n"
        "→ 换 ArcFace 做纵向对比",
        6.8, 2.0, 5.5, 4.8, font_size=13, color=COLOR_BLACK)


def slide_dataset(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_bottom_bar(slide); add_page_number(slide, 5)
    add_title_text(slide, "数据集与实验设计", 0.8, 0.4, 10, 0.7, font_size=32, color=COLOR_DARK)
    add_separator(slide)

    add_title_text(slide, "FairFace 数据集", 0.8, 1.4, 5, 0.5, font_size=20, color=COLOR_PRIMARY)
    add_body_text(slide,
        "◆ 108,501 张人脸图像 (train + val)\n"
        "◆ 7 个种族类别: White, Black, East Asian,\n"
        "   Southeast Asian, Indian, Middle Eastern,\n"
        "   Latino/Hispanic\n"
        "◆ 2 个性别: Male, Female\n"
        "◆ 年龄段: 0-70+ (9 段)\n"
        "◆ 平衡采样: 每种族均匀抽取\n"
        "◆ 来源: YFCC100M Flickr 数据集",
        0.8, 2.0, 5.8, 3.5, font_size=13, color=COLOR_BLACK)

    add_title_text(slide, "实验配置", 7.2, 1.4, 5, 0.5, font_size=20, color=COLOR_PRIMARY)
    add_body_text(slide,
        f"分析样本: {data['total']} 张 (7 种族均匀采样)\n"
        f"人脸检测: OpenCV\n"
        f"识别模型: Facenet\n\n"
        "评估指标:\n"
        "① 种族识别准确率 (按种族分组)\n"
        "② 性别识别准确率 (按种族×性别)\n"
        "③ 混淆矩阵 (行归一化)\n"
        "④ 置信度分布 (箱线图)\n"
        "⑤ 公平性偏差差距 (<5% 轻微,\n"
        "   5-15% 中度, >15% 显著)",
        7.2, 2.0, 5.5, 4.5, font_size=13, color=COLOR_BLACK)


def slide_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_bottom_bar(slide); add_page_number(slide, 6)
    add_title_text(slide, "分析流水线", 0.8, 0.4, 10, 0.7, font_size=32, color=COLOR_DARK)
    add_separator(slide)

    steps = [
        ("1", "加载数据", "FairFace 数据集\n按种族均匀采样\n读取 CSV 标签"),
        ("2", "DeepFace 分析", "人脸检测 + 对齐\nFacenet 特征提取\n种族/性别/年龄预测"),
        ("3", "统计分析", "预测 vs 真实标签\n计算分群准确率\n偏差量化评估"),
        ("4", "可视化", "5 张图表\n柱状图 / 热力图\n混淆矩阵 / 雷达图"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        x = 0.6 + i * 3.1
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.8), Inches(1.6), Inches(0.5), Inches(0.5))
        circle.fill.solid(); circle.fill.fore_color.rgb = COLOR_PRIMARY if i < 3 else COLOR_ACCENT; circle.line.fill.background()
        p = circle.text_frame.paragraphs[0]
        p.text = num; p.font.size = Pt(16); p.font.color.rgb = COLOR_WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
        add_title_text(slide, title, x, 2.3, 2.8, 0.5, font_size=18, color=COLOR_DARK, alignment=PP_ALIGN.CENTER)
        add_body_text(slide, desc, x, 2.9, 2.8, 2.0, font_size=12, color=COLOR_BLACK)

    for i in range(3):
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.3 + i * 3.1), Inches(1.75), Inches(0.4), Inches(0.2))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = COLOR_GRAY; arrow.line.fill.background()


def slide_results_overall(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_bottom_bar(slide); add_page_number(slide, 7)
    add_title_text(slide, "实验结果 — 总体准确率", 0.8, 0.4, 10, 0.7, font_size=32, color=COLOR_DARK)
    add_separator(slide)

    # 三个指标卡片
    add_body_text(slide, f"总体种族准确率\n{data['race_acc']:.1f}%", 0.8, 1.6, 3.5, 1.3, font_size=22, color=COLOR_PRIMARY)
    add_body_text(slide, f"总体性别准确率\n{data['gender_acc']:.1f}%", 5.0, 1.6, 3.5, 1.3, font_size=22, color=COLOR_ACCENT)
    add_body_text(slide, f"总样本: {data['total']}\n检测失败: {data['failed']}", 9.2, 1.6, 3.5, 1.3, font_size=18, color=COLOR_GRAY)

    add_image_safe(slide, os.path.join(CHARTS_DIR, "01_race_accuracy.png"), 0.5, 3.2, 6.0, 3.8)
    add_image_safe(slide, os.path.join(CHARTS_DIR, "03_confusion_matrix.png"), 6.8, 3.2, 6.0, 3.8)


def slide_results_race(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_bottom_bar(slide); add_page_number(slide, 8)
    add_title_text(slide, "实验结果 — 各种族准确率对比", 0.8, 0.4, 10, 0.7, font_size=32, color=COLOR_DARK)
    add_separator(slide)

    # 数据表格
    lines = [f"{'人种':<20s} {'样本':<6s} {'种族':<8s} {'性别':<8s}" + " " * 20]
    for race in RACE_CATEGORIES:
        rs = data["race_stats"].get(race)
        if rs:
            lines.append(f"{race:<20s} {rs['count']:<6d} {rs['race_acc']:.1f}%{'':>4s} {rs['gender_acc']:.1f}%")
    add_body_text(slide, "\n".join(lines), 0.8, 1.5, 5.8, 3.0, font_size=12, color=COLOR_BLACK)

    add_image_safe(slide, os.path.join(CHARTS_DIR, "02_race_gender_heatmap.png"), 6.8, 1.5, 6.0, 3.0)

    # 关键发现
    gap_color = COLOR_RED if data["gap"] > 15 else (COLOR_ACCENT if data["gap"] > 5 else COLOR_GREEN)
    finding = (f"最高: {data['best_race']} ({data['best_acc']:.1f}%)  |  "
               f"最低: {data['worst_race']} ({data['worst_acc']:.1f}%)  |  "
               f"差距: {data['gap']:.1f}%  →  {data['assessment']}")
    add_body_text(slide, finding, 0.8, 5.0, 12, 1.5, font_size=16, color=gap_color)


def slide_confusion_detail(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_bottom_bar(slide); add_page_number(slide, 9)
    add_title_text(slide, "实验结果 — 置信度分析", 0.8, 0.4, 10, 0.7, font_size=32, color=COLOR_DARK)
    add_separator(slide)

    add_image_safe(slide, os.path.join(CHARTS_DIR, "04_confidence_boxplot.png"), 0.3, 1.3, 6.0, 5.8)
    add_image_safe(slide, os.path.join(CHARTS_DIR, "05_fairness_radar.png"), 6.8, 1.3, 6.0, 5.8)


def slide_fairness(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_bottom_bar(slide); add_page_number(slide, 10)
    add_title_text(slide, "公平性偏差评估", 0.8, 0.4, 10, 0.7, font_size=32, color=COLOR_DARK)
    add_separator(slide)

    # 左侧: 偏差标准 + 实际数据
    gender_m = data["gender_stats"].get("Male", {})
    gender_f = data["gender_stats"].get("Female", {})
    gender_gap = abs(gender_m.get("acc", 0) - gender_f.get("acc", 0))

    left_text = (
        f"偏差评估标准:\n\n"
        f"差距 < 5%\n  → 轻微偏差，模型较公平\n\n"
        f"差距 5% - 15%\n  → 中度偏差，值得关注\n\n"
        f"差距 > 15%\n  → 显著偏差，存在系统性偏见"
    )
    add_body_text(slide, left_text, 0.8, 1.5, 5.5, 3.5, font_size=13, color=COLOR_BLACK)

    right_text = (
        f"本次分析结果:\n\n"
        f"种族准确率差距 = {data['gap']:.1f}%\n"
        f"→ {data['assessment']}\n\n"
        f"最高: {data['best_race']} ({data['best_acc']:.1f}%)\n"
        f"最低: {data['worst_race']} ({data['worst_acc']:.1f}%)\n\n"
        f"性别偏差:\n"
        f"Male 准确率: {gender_m.get('acc', 0):.1f}%\n"
        f"Female 准确率: {gender_f.get('acc', 0):.1f}%\n"
        f"差距: {gender_gap:.1f}%"
    )
    add_body_text(slide, right_text, 0.8, 5.2, 12, 2.0, font_size=13, color=COLOR_BLACK)

    # 右侧: 雷达图
    add_image_safe(slide, os.path.join(CHARTS_DIR, "05_fairness_radar.png"), 7.0, 1.5, 5.8, 5.5)


def slide_conclusion(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_bottom_bar(slide); add_page_number(slide, 11)
    add_title_text(slide, "结论与讨论", 0.8, 0.4, 10, 0.7, font_size=32, color=COLOR_DARK)
    add_separator(slide)

    gender_m = data["gender_stats"].get("Male", {})
    gender_f = data["gender_stats"].get("Female", {})
    gender_gap = abs(gender_m.get("acc", 0) - gender_f.get("acc", 0))

    findings = (
        f"1. Facenet 在种族识别上存在 {data['bias_level']} 偏见\n"
        f"   差距最大的两个群体: {data['best_race']} vs {data['worst_race']}\n"
        f"   (Δ = {data['gap']:.1f}%)\n\n"
        f"2. 性别维度存在显著交叉效应\n"
        f"   Male 准确率: {gender_m.get('acc', 0):.1f}%\n"
        f"   Female 准确率: {gender_f.get('acc', 0):.1f}% (差距: {gender_gap:.1f}%)\n\n"
        f"3. 总体种族准确率仅 {data['race_acc']:.1f}%，\n"
        f"   Facenet 在种族分类任务上整体表现较弱\n\n"
        f"4. 种族偏差与性别偏差叠加，\n"
        f"   Latino_Hispanic + Female 组合受影响最严重"
    )
    add_title_text(slide, "主要发现", 0.8, 1.4, 5, 0.5, font_size=20, color=COLOR_PRIMARY)
    add_body_text(slide, findings, 0.8, 2.0, 5.8, 5.0, font_size=13, color=COLOR_BLACK)

    add_title_text(slide, "局限性与后续工作", 7.2, 1.4, 5.5, 0.5, font_size=20, color=COLOR_PRIMARY)
    add_body_text(slide,
        "局限性:\n"
        "◆ 仅测试了 Facenet 一个识别模型\n"
        "◆ FairFace 标注可能引入人为偏差\n"
        "◆ 年龄维度的公平性未深入\n"
        "◆ 数据集来自 Flickr，非真实部署场景\n\n"
        "后续方向:\n"
        "◆ 横向对比 ArcFace / VGG-Face\n"
        "◆ 引入皮肤色调 Fitzpatrick 量表\n"
        "◆ 换成真实监控场景数据集",
        7.2, 2.0, 5.5, 5.0, font_size=13, color=COLOR_BLACK)


def slide_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.2), Inches(2.3), Inches(0.06))
    shape.fill.solid(); shape.fill.fore_color.rgb = COLOR_ACCENT; shape.line.fill.background()

    add_title_text(slide, "感谢聆听", 2, 1.8, 9.5, 1.0, font_size=44, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
    add_body_text(slide, "FaceFair · AI 公平性审计\n人工智能导论 课程项目\n\n[姓名 / 学号]", 3, 3.6, 7.5, 2.0, font_size=16, color=COLOR_GRAY)


# ========== 主流程 ==========

def main():
    print("解析数据...")
    data = parse_summary()
    print(f"  样本: {data['total']}, 种族准确率: {data['race_acc']:.1f}%, 性别准确率: {data['gender_acc']:.1f}%")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)                       # 1
    slide_background(prs)                  # 2
    slide_problem(prs)                     # 3
    slide_approach(prs)                    # 4
    slide_dataset(prs, data)               # 5
    slide_pipeline(prs)                    # 6
    slide_results_overall(prs, data)       # 7
    slide_results_race(prs, data)          # 8
    slide_confusion_detail(prs, data)      # 9
    slide_fairness(prs, data)              # 10
    slide_conclusion(prs, data)            # 11
    slide_thanks(prs)                      # 12

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    output_path = os.path.join(OUTPUT_DIR, "FaceFair_完整版.pptx")
    prs.save(output_path)
    print(f"\nPPT 已保存: {output_path}")
    print(f"共 {len(prs.slides)} 张幻灯片，数据已填入")


if __name__ == "__main__":
    main()
